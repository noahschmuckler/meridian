#!/usr/bin/env python3
"""Generate a PowerPoint presentation from a Meridian module JSON file."""

import json
import os
import re
import sys
from html.parser import HTMLParser

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Meridian Colors ──
GREEN = RGBColor(0x1A, 0x5C, 0x3A)
GREEN_LIGHT = RGBColor(0xEA, 0xF2, 0xED)
GREEN_MID = RGBColor(0x2D, 0x7A, 0x52)
BROWN = RGBColor(0x7A, 0x3B, 0x1E)
BROWN_LIGHT = RGBColor(0xF5, 0xED, 0xE8)
TEXT_DARK = RGBColor(0x1C, 0x1A, 0x16)
TEXT_MUTED = RGBColor(0x6B, 0x65, 0x60)
BG_WARM = RGBColor(0xF4, 0xF1, 0xEC)
SURFACE = RGBColor(0xFD, 0xFC, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xD6, 0xD0, 0xC4)

FONT = 'Segoe UI'
FONT_MONO = 'Segoe UI'  # Closest match; no mono needed in slides

# Slide dimensions: standard widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ── HTML to plain text with basic formatting info ──
class HTMLToRuns(HTMLParser):
    """Parse FAQ answer HTML into a list of text runs with formatting."""

    def __init__(self):
        super().__init__()
        self.runs = []  # list of {text, bold, italic, pill, pill_type}
        self._bold = False
        self._italic = False
        self._pill = False
        self._pill_type = None
        self._in_p = False
        self._p_count = 0

    def handle_starttag(self, tag, attrs):
        attrs_dict = dict(attrs)
        if tag == 'strong':
            self._bold = True
        elif tag == 'em':
            self._italic = True
        elif tag == 'p':
            if self._p_count > 0:
                self.runs.append({'text': '\n', 'bold': False, 'italic': False, 'pill': False})
            self._p_count += 1
            self._in_p = True
        elif tag == 'span':
            cls = attrs_dict.get('class', '')
            if 'pill' in cls:
                self._pill = True
                if 'warn' in cls:
                    self._pill_type = 'warn'
                elif 'green' in cls:
                    self._pill_type = 'green'
                else:
                    self._pill_type = 'neutral'

    def handle_endtag(self, tag):
        if tag == 'strong':
            self._bold = False
        elif tag == 'em':
            self._italic = False
        elif tag == 'p':
            self._in_p = False
        elif tag == 'span':
            if self._pill:
                self._pill = False
                self._pill_type = None

    def handle_data(self, data):
        if data.strip() or data == '\n':
            self.runs.append({
                'text': data,
                'bold': self._bold,
                'italic': self._italic,
                'pill': self._pill,
                'pill_type': self._pill_type,
            })


def parse_answer_html(html):
    """Convert answer HTML to a list of formatted text runs."""
    parser = HTMLToRuns()
    parser.feed(html)
    return parser.runs


def html_to_plain(html):
    """Quick HTML to plain text."""
    text = re.sub(r'<br\s*/?>', '\n', html)
    text = re.sub(r'<[^>]+>', '', text)
    text = re.sub(r'&amp;', '&', text)
    text = re.sub(r'&lt;', '<', text)
    text = re.sub(r'&gt;', '>', text)
    text = re.sub(r'&quot;', '"', text)
    text = re.sub(r'&#39;', "'", text)
    text = re.sub(r'&nbsp;', ' ', text)
    return text.strip()


# ── Slide Helpers ──

def set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height):
    """Add a textbox and return its text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def add_run(paragraph, text, size=14, color=TEXT_DARK, bold=False, italic=False, font=FONT):
    """Add a formatted run to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return run


def add_paragraph(tf, text='', size=14, color=TEXT_DARK, bold=False, italic=False,
                  alignment=PP_ALIGN.LEFT, space_after=Pt(4), space_before=Pt(0)):
    """Add a new paragraph to a text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    if text:
        add_run(p, text, size=size, color=color, bold=bold, italic=italic)
    return p


# ── Slide Builders ──

def build_title_slide(prs, mod):
    """Slide 1: Green background, module title, subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, GREEN)

    # Eyebrow
    tf = add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5))
    add_paragraph(tf, 'MERIDIAN', size=12, color=WHITE, bold=False)
    tf.paragraphs[0].font.name = FONT

    # Title
    parts = mod['default_title'].split(' \u2014 ', 1)
    title = parts[0] if len(parts) == 2 else mod['default_title']
    subtitle = parts[1] if len(parts) == 2 else ''

    tf = add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5))
    add_paragraph(tf, title, size=36, color=WHITE, bold=True)

    if subtitle:
        tf2 = add_textbox(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8))
        add_paragraph(tf2, subtitle, size=20, color=RGBColor(0xCC, 0xDD, 0xCC), italic=True)

    # Intro text
    if mod.get('landing_intro'):
        tf3 = add_textbox(slide, Inches(1), Inches(4.8), Inches(9), Inches(1.5))
        add_paragraph(tf3, mod['landing_intro'], size=14, color=RGBColor(0xCC, 0xDD, 0xCC), italic=True)

    # Bottom bar
    add_shape(slide, Inches(1), Inches(6.8), Inches(11.333), Pt(2), fill_color=RGBColor(0xCC, 0xDD, 0xCC))


def build_checklist_slide(prs, mod, faq_ref_map):
    """Slide 2: Checklist items + green zone."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)

    # Header bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.8), fill_color=GREEN)
    tf = add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.5))
    add_paragraph(tf, mod['checklist_section_label'].upper(), size=11, color=WHITE, bold=True)

    # Checklist items
    y = Inches(1.2)
    for item in mod['checklist']:
        ref_label = faq_ref_map.get(item['faq_ref'], '')

        # Checkbox shape
        cb = add_shape(slide, Inches(0.8), y + Pt(3), Inches(0.22), Inches(0.22),
                       line_color=BORDER, line_width=Pt(1.5))

        # Item text
        tf = add_textbox(slide, Inches(1.2), y, Inches(9.5), Inches(0.6))
        p = add_paragraph(tf, '', size=15, color=TEXT_DARK)
        add_run(p, item['statement'], size=15, color=TEXT_DARK)
        if ref_label:
            add_run(p, '  ' + ref_label, size=10, color=GREEN, bold=True)

        y += Inches(0.65)

    # Green zone box
    y += Inches(0.3)
    gz = mod.get('green_zone')
    if gz:
        # Left accent bar
        add_shape(slide, Inches(0.7), y, Pt(4), Inches(1.6), fill_color=GREEN)
        # Background
        add_shape(slide, Inches(0.75), y, Inches(10.5), Inches(1.6), fill_color=GREEN_LIGHT)

        # Zone label
        tf = add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(10), Inches(0.3))
        add_paragraph(tf, gz['zone_label'].upper(), size=10, color=GREEN, bold=True)

        # Narrative
        narrative = html_to_plain(gz['narrative_html'])
        tf2 = add_textbox(slide, Inches(1.0), y + Inches(0.45), Inches(10), Inches(0.6))
        add_paragraph(tf2, narrative, size=13, color=TEXT_DARK)

        # SmartPhrase tag
        tf3 = add_textbox(slide, Inches(1.0), y + Inches(1.1), Inches(4), Inches(0.35))
        p = add_paragraph(tf3, '', size=11, color=GREEN)
        add_run(p, 'SmartPhrase: ' + gz['smartphrase'], size=11, color=GREEN, bold=False, font=FONT)


def build_escalation_slide(prs, mod, faq_ref_map):
    """Slide 3: Escalation indicators."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)

    # Header bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.8), fill_color=BROWN)
    tf = add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.5))
    add_paragraph(tf, mod['escalation_section_label'].upper(), size=11, color=WHITE, bold=True)

    # Escalation items
    y = Inches(1.2)
    item_height = Inches(0.55) if len(mod['escalation']) <= 6 else Inches(0.48)
    font_size = 14 if len(mod['escalation']) <= 6 else 13

    for item in mod['escalation']:
        ref_label = faq_ref_map.get(item['faq_ref'], '')

        # Dot
        dot = add_shape(slide, Inches(0.8), y + Pt(5), Inches(0.1), Inches(0.1), fill_color=BROWN)

        # Item text
        tf = add_textbox(slide, Inches(1.1), y, Inches(10), Inches(0.5))
        p = add_paragraph(tf, '', size=font_size, color=TEXT_DARK)
        add_run(p, item['statement'], size=font_size, color=TEXT_DARK)
        if ref_label:
            add_run(p, '  ' + ref_label, size=10, color=GREEN, bold=True)

        y += item_height

    # Context strip if present
    if mod.get('context_strip'):
        y += Inches(0.3)
        add_shape(slide, Inches(0.7), y, Inches(10.5), Inches(1.0),
                  fill_color=WHITE, line_color=BORDER, line_width=Pt(1))
        tf = add_textbox(slide, Inches(0.9), y + Inches(0.08), Inches(10), Inches(0.2))
        add_paragraph(tf, mod['context_strip']['label'].upper(), size=9, color=TEXT_MUTED, bold=True)
        tf2 = add_textbox(slide, Inches(0.9), y + Inches(0.3), Inches(10), Inches(0.6))
        add_paragraph(tf2, mod['context_strip']['text'], size=12, color=TEXT_MUTED, italic=True)


def build_faq_slide(prs, mod, faq, ref_label):
    """One slide per FAQ entry with all Q&A pairs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)

    # Header bar
    add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.0), fill_color=GREEN)

    # Ref badge
    tf_badge = add_textbox(slide, Inches(0.6), Inches(0.15), Inches(0.8), Inches(0.35))
    p = add_paragraph(tf_badge, ref_label, size=14, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    badge_shape = add_shape(slide, Inches(0.55), Inches(0.13), Inches(0.65), Inches(0.35),
                            fill_color=WHITE, line_color=None)
    # Move badge behind text by reordering (badge first, then text on top)
    # Actually just overlay — the textbox is on top by default

    # Topic
    tf_topic = add_textbox(slide, Inches(1.4), Inches(0.12), Inches(10), Inches(0.3))
    add_paragraph(tf_topic, faq['topic'].upper(), size=10, color=RGBColor(0xCC, 0xDD, 0xCC))

    # Title
    tf_title = add_textbox(slide, Inches(1.4), Inches(0.42), Inches(10), Inches(0.5))
    add_paragraph(tf_title, faq['title'], size=18, color=WHITE, bold=True)

    # Q&A pairs
    y = Inches(1.3)
    max_y = Inches(6.8)

    for qa in faq['items']:
        if y > max_y:
            # Overflow — start new slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide, SURFACE)
            add_shape(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.6), fill_color=GREEN)
            tf_cont = add_textbox(slide, Inches(0.6), Inches(0.1), Inches(10), Inches(0.4))
            add_paragraph(tf_cont, ref_label + ' — ' + faq['title'] + ' (continued)',
                          size=13, color=WHITE, bold=True)
            y = Inches(0.9)

        # Question
        tf_q = add_textbox(slide, Inches(0.6), y, Inches(11.5), Inches(0.5))
        add_paragraph(tf_q, qa['question'], size=14, color=TEXT_DARK, bold=True,
                      space_after=Pt(2))

        y += Inches(0.45)

        # Answer — render with formatting
        runs = parse_answer_html(qa['answer_html'])
        tf_a = add_textbox(slide, Inches(0.8), y, Inches(11), Inches(3))

        # Left accent border
        add_shape(slide, Inches(0.65), y, Pt(2.5), Inches(0.01), fill_color=GREEN_LIGHT)

        p = None
        for run_info in runs:
            text = run_info['text']
            if text == '\n':
                p = None
                continue
            if p is None:
                p = add_paragraph(tf_a, '', size=12, space_after=Pt(6))

            if run_info['pill']:
                # Render pills as bracketed bold text
                pill_color = BROWN if run_info['pill_type'] == 'warn' else GREEN if run_info['pill_type'] == 'green' else TEXT_MUTED
                add_run(p, '[' + text + ']', size=11, color=pill_color, bold=True)
            else:
                add_run(p, text, size=12,
                        color=TEXT_DARK if run_info['bold'] else TEXT_MUTED,
                        bold=run_info['bold'],
                        italic=run_info['italic'])

        # Estimate height used by answer (rough: ~20pt per line, ~80 chars per line)
        total_text = html_to_plain(qa['answer_html'])
        est_lines = max(1, len(total_text) // 90 + total_text.count('\n') + 1)
        y += Inches(0.25 * est_lines + 0.15)

        # Separator
        if y < max_y:
            add_shape(slide, Inches(0.8), y, Inches(10.5), Pt(0.75), fill_color=BORDER)
            y += Inches(0.15)


def build_footer_slide(prs, mod):
    """Final slide with footer note."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, GREEN)

    tf = add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1))
    add_paragraph(tf, 'MERIDIAN', size=12, color=RGBColor(0xCC, 0xDD, 0xCC))

    parts = mod['default_title'].split(' \u2014 ', 1)
    title = parts[0] if len(parts) == 2 else mod['default_title']
    tf2 = add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(1))
    add_paragraph(tf2, title, size=28, color=WHITE, bold=True)

    if mod.get('footer_note'):
        tf3 = add_textbox(slide, Inches(1), Inches(4.5), Inches(9), Inches(1.5))
        add_paragraph(tf3, mod['footer_note'], size=13, color=RGBColor(0xCC, 0xDD, 0xCC), italic=True)


# ── Main ──

def generate_pptx(module_path, output_path=None):
    """Generate a PPTX from a module JSON file."""
    with open(module_path, 'r', encoding='utf-8') as f:
        mod = json.load(f)

    if output_path is None:
        output_path = os.path.splitext(module_path)[0] + '.pptx'

    # Build FAQ reference map
    faq_ref_map = {}
    for i, faq in enumerate(mod['faqs']):
        faq_ref_map[faq['faq_id']] = 'F' + str(i + 1)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build slides
    build_title_slide(prs, mod)
    build_checklist_slide(prs, mod, faq_ref_map)
    build_escalation_slide(prs, mod, faq_ref_map)

    for i, faq in enumerate(mod['faqs']):
        ref_label = 'F' + str(i + 1)
        build_faq_slide(prs, mod, faq, ref_label)

    build_footer_slide(prs, mod)

    prs.save(output_path)
    return output_path


def main():
    if len(sys.argv) < 2:
        # Default: generate all modules
        modules_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'modules')
        index_path = os.path.join(modules_dir, 'index.json')
        with open(index_path) as f:
            index = json.load(f)

        output_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'presentations')
        os.makedirs(output_dir, exist_ok=True)

        for entry in index['modules']:
            mod_path = os.path.join(modules_dir, entry['file'])
            out_path = os.path.join(output_dir, entry['module_id'] + '.pptx')
            print(f"Generating {entry['module_id']}...")
            generate_pptx(mod_path, out_path)
            print(f"  -> {out_path}")

        print(f"\nDone. {len(index['modules'])} presentations generated.")
    else:
        module_path = sys.argv[1]
        output_path = sys.argv[2] if len(sys.argv) > 2 else None
        result = generate_pptx(module_path, output_path)
        print(f"Generated: {result}")


if __name__ == '__main__':
    main()
